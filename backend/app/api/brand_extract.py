from __future__ import annotations

import asyncio
import base64
import copy
import hashlib
import io
import json as _json
import logging
import re
import uuid
from typing import Any

import httpx
from fastapi import APIRouter, Depends, File, Form, HTTPException, UploadFile, status

from ..auth import AuthUser, get_library_user
from ..config import settings

logger = logging.getLogger("prezo.brand_extract")

router = APIRouter(prefix="/library/poll-game/brand-profiles", tags=["brand-extract"])

EXTRACT_TIMEOUT_SECONDS = 180.0
EXTRACT_MAX_FILE_SIZE = 50 * 1024 * 1024   # 50 MB hard cap
INLINE_MAX_FILE_SIZE = 14 * 1024 * 1024    # >14 MB → use File API (base64 would exceed 20 MB limit)

# Image extraction settings
IMAGE_MIN_DIMENSION = 40      # skip tiny images (icons, bullets)
IMAGE_MAX_CANDIDATES = 10     # return at most N candidate images
IMAGE_MAX_DATA_URL_BYTES = 500_000  # skip images whose base64 exceeds ~500 KB

SUPPORTED_IMAGE_TYPES = {"image/png", "image/jpeg", "image/gif", "image/webp"}
SUPPORTED_UPLOAD_TYPES = SUPPORTED_IMAGE_TYPES | {
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

# ---------------------------------------------------------------------------
# Pass 1 — Visual Identity
# Focused on: colours, typography, patterns, shapes, asset styles
# ---------------------------------------------------------------------------

PASS1_SYSTEM = (
    "You are a meticulous brand identity analyst. Given the uploaded content "
    "(brand guidelines PDF, presentation slides, logo image, or website screenshot), "
    "extract the brand's VISUAL IDENTITY SYSTEM.\n\n"
    "Focus ONLY on visual/design elements:\n"
    "- Colours (primary, secondary, accent) with hex values, names, and usage\n"
    "- Gradients with definitions and usage\n"
    "- Fonts with family names, weights, and usage roles\n"
    "- Typography hierarchy (heading/body/caption relationships)\n"
    "- Logo description, variations, rules, and colours\n"
    "- Patterns, textures, scribbles, and decorative elements\n"
    "- Brand shapes and geometric elements\n"
    "- Background treatment styles\n"
    "- Iconography style and rules\n"
    "- Illustration style and approach\n"
    "- Photography style and treatment\n"
    "- Spacing and layout grid rules\n"
    "- Animation/motion guidelines\n\n"
    "Be EXHAUSTIVE. Extract every colour swatch with its hex, RGB, CMYK, and "
    "Pantone values if provided. List every font weight. Describe every pattern. "
    "For a 100+ page brand guide, your response should be very detailed.\n\n"
    "Return a JSON object matching the required schema."
)

PASS1_SCHEMA: dict[str, Any] = {
    "type": "object",
    "properties": {
        "primary_colors": {
            "type": "array",
            "items": {"type": "string"},
            "description": "Hex colour strings with usage notes, e.g. '#1A2B3C – Navy, for headlines'",
        },
        "secondary_colors": {
            "type": "array",
            "items": {"type": "string"},
        },
        "accent_colors": {
            "type": "array",
            "items": {"type": "string"},
        },
        "gradient_styles": {
            "type": "array",
            "items": {"type": "string"},
            "description": "Gradient definitions with usage",
        },
        "fonts": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "family": {"type": "string"},
                    "weights": {"type": "array", "items": {"type": "string"}},
                    "usage": {"type": "string"},
                },
                "required": ["family"],
                "additionalProperties": False,
            },
        },
        "typography_hierarchy": {"type": "string"},
        "logo_description": {"type": "string"},
        "logo_colors": {
            "type": "array",
            "items": {"type": "string"},
        },
        "patterns_and_textures": {"type": "string"},
        "brand_shapes": {"type": "string"},
        "background_styles": {"type": "string"},
        "iconography_style": {"type": "string"},
        "illustration_style": {"type": "string"},
        "photography_style": {"type": "string"},
        "spacing_and_layout": {"type": "string"},
        "animation_motion": {"type": "string"},
    },
    "required": [
        "primary_colors", "secondary_colors", "accent_colors",
        "gradient_styles", "fonts", "typography_hierarchy",
        "logo_description", "logo_colors",
        "patterns_and_textures", "brand_shapes", "background_styles",
        "iconography_style", "illustration_style", "photography_style",
        "spacing_and_layout", "animation_motion",
    ],
    "additionalProperties": False,
}

# ---------------------------------------------------------------------------
# Pass 2 — Brand Voice & Strategy
# Focused on: tone, messaging, principles, dos/don'ts, visual style philosophy
# ---------------------------------------------------------------------------

PASS2_SYSTEM = (
    "You are a meticulous brand identity analyst. Given the uploaded content "
    "(brand guidelines PDF, presentation slides, logo image, or website screenshot), "
    "extract the brand's VOICE, STRATEGY, AND DESIGN PHILOSOPHY.\n\n"
    "Focus ONLY on:\n"
    "- Overall visual style philosophy and aesthetic direction\n"
    "- Tone of voice (personality, communication style, formality level)\n"
    "- Key brand principles and values\n"
    "- Messaging framework (taglines, slogans, key messages, value propositions)\n"
    "- Do's and Don'ts (explicit rules from the guidelines)\n\n"
    "Be EXHAUSTIVE. Extract every principle, every do/don't, every messaging "
    "example. Quote directly from the guidelines where possible. For a 100+ page "
    "brand guide, your response should be very detailed.\n\n"
    "Return a JSON object matching the required schema."
)

PASS2_SCHEMA: dict[str, Any] = {
    "type": "object",
    "properties": {
        "visual_style": {
            "type": "string",
            "description": "Detailed summary of the brand's visual tone, aesthetic, and design philosophy",
        },
        "tone_of_voice": {
            "type": "string",
            "description": "Communication style, personality traits, formality level",
        },
        "key_principles": {
            "type": "array",
            "items": {"type": "string"},
            "description": "ALL brand guidelines principles found",
        },
        "messaging_framework": {
            "type": "array",
            "items": {"type": "string"},
            "description": "Key messages, taglines, slogans, value propositions",
        },
        "dos_and_donts": {
            "type": "array",
            "items": {"type": "string"},
            "description": "Explicit do's and don'ts from the guidelines",
        },
    },
    "required": [
        "visual_style", "tone_of_voice", "key_principles",
        "messaging_framework", "dos_and_donts",
    ],
    "additionalProperties": False,
}


# ---------------------------------------------------------------------------
# Pass 3 — UI identity (structured brand name, 6 color roles, typography)
# Used for the host Brand identity editor; stored under guidelines["ui_identity"].
# ---------------------------------------------------------------------------

PASS3_SYSTEM = (
    "You extract a compact UI IDENTITY block for a presentation product (slides and on-screen UI).\n\n"
    "Return exactly 6 color roles, each with a clear job on a typical slide:\n"
    "1) Slide background 2) Headline / title text 3) Accent & highlights (buttons, icons, callouts)\n"
    "4) Card & panel fill 5) Body / paragraph text 6) Borders & dividers\n\n"
    "For each color role provide:\n"
    "- role: short title (e.g. 'Slide Background')\n"
    "- usage: one line describing where it is used\n"
    "- hex: primary sRGB hex like #RRGGBB (prefer colours from the brand document)\n"
    "- hierarchy_rank: integer 1–6 where 1 = most visually dominant / largest area on a typical slide, "
    "6 = least dominant (thin lines / rare accents)\n"
    "- surface: one of background, foreground, accent, border, fill, neutral — best match for that swatch\n\n"
    "brand_name: the official brand or product name from the document (or a reasonable inferred name).\n\n"
    "typography:\n"
    "- heading_1.family, heading_2.family, body.family — font family names used for main title, "
    "subhead/section titles, and body copy (infer from guidelines if not explicit).\n\n"
    "Be concise. Use real hex values. Do not invent brand names that contradict the document.\n\n"
    "tone_calibration — four axes for voice, each an integer 0–100 (0 = left pole, 100 = right pole, 50 = balanced):\n"
    "- serious_playful: 0 = serious, 100 = playful\n"
    "- formal_casual: 0 = formal, 100 = casual\n"
    "- respectful_irreverent: 0 = respectful, 100 = irreverent\n"
    "- matter_of_fact_enthusiastic: 0 = matter-of-fact, 100 = enthusiastic\n"
    "Infer from tone of voice, messaging, audience, and style cues in the document.\n\n"
    "visual_style — write like a concise brand style guide (plain sentences; NO markdown, NO bullet lists):\n"
    "- visual_mood_aesthetic: One or two short sentences summarizing overall look, mood, and aesthetic "
    "(e.g. traits strung together in flowing prose: modern, vibrant, approachable…).\n"
    "- style_guidelines: One or two short paragraphs on how layouts, typography, emphasis, and composition "
    "should feel on slides — concrete and specific to the document.\n"
    "- design_elements: four STRING fields (patterns_textures, icon_style, image_treatment, decorative_elements). "
    "Each must be one or two sentences of descriptive prose (how that aspect should look), not comma-separated "
    "keyword lists or labels only. Quote or closely paraphrase the guidelines when possible.\n"
    "If the document is silent on a field, use an empty string. Do not invent contradictory details."
)

_VISUAL_STYLE_TEXT_MAX = 8000

_VISUAL_STYLE_FIELD: dict[str, Any] = {
    "type": "string",
    "maxLength": _VISUAL_STYLE_TEXT_MAX,
}

PASS3_SCHEMA: dict[str, Any] = {
    "type": "object",
    "properties": {
        "brand_name": {"type": "string"},
        "color_roles": {
            "type": "array",
            "minItems": 6,
            "maxItems": 6,
            "items": {
                "type": "object",
                "properties": {
                    "role": {"type": "string"},
                    "usage": {"type": "string"},
                    "hex": {"type": "string"},
                    "hierarchy_rank": {"type": "integer"},
                    "surface": {"type": "string"},
                },
                "required": ["role", "usage", "hex", "hierarchy_rank", "surface"],
                "additionalProperties": False,
            },
        },
        "typography": {
            "type": "object",
            "properties": {
                "heading_1": {
                    "type": "object",
                    "properties": {"family": {"type": "string"}},
                    "required": ["family"],
                    "additionalProperties": False,
                },
                "heading_2": {
                    "type": "object",
                    "properties": {"family": {"type": "string"}},
                    "required": ["family"],
                    "additionalProperties": False,
                },
                "body": {
                    "type": "object",
                    "properties": {"family": {"type": "string"}},
                    "required": ["family"],
                    "additionalProperties": False,
                },
            },
            "required": ["heading_1", "heading_2", "body"],
            "additionalProperties": False,
        },
        "tone_calibration": {
            "type": "object",
            "properties": {
                "serious_playful": {"type": "integer", "minimum": 0, "maximum": 100},
                "formal_casual": {"type": "integer", "minimum": 0, "maximum": 100},
                "respectful_irreverent": {"type": "integer", "minimum": 0, "maximum": 100},
                "matter_of_fact_enthusiastic": {"type": "integer", "minimum": 0, "maximum": 100},
            },
            "required": [
                "serious_playful",
                "formal_casual",
                "respectful_irreverent",
                "matter_of_fact_enthusiastic",
            ],
            "additionalProperties": False,
        },
        "visual_style": {
            "type": "object",
            "properties": {
                "visual_mood_aesthetic": _VISUAL_STYLE_FIELD,
                "style_guidelines": _VISUAL_STYLE_FIELD,
                "design_elements": {
                    "type": "object",
                    "properties": {
                        "patterns_textures": _VISUAL_STYLE_FIELD,
                        "icon_style": _VISUAL_STYLE_FIELD,
                        "image_treatment": _VISUAL_STYLE_FIELD,
                        "decorative_elements": _VISUAL_STYLE_FIELD,
                    },
                    "required": [
                        "patterns_textures",
                        "icon_style",
                        "image_treatment",
                        "decorative_elements",
                    ],
                    "additionalProperties": False,
                },
            },
            "required": ["visual_mood_aesthetic", "style_guidelines", "design_elements"],
            "additionalProperties": False,
        },
    },
    "required": ["brand_name", "color_roles", "typography", "tone_calibration", "visual_style"],
    "additionalProperties": False,
}

_DEFAULT_COLOR_ROLES: list[tuple[str, str, str]] = [
    ("Slide Background", "Main slide background", "background"),
    ("Headline Text", "Titles and headings", "foreground"),
    ("Accent & Highlights", "Buttons, icons, callouts", "accent"),
    ("Card & Panel Fill", "Cards, boxes, containers", "fill"),
    ("Body Text", "Paragraphs and captions", "foreground"),
    ("Borders & Dividers", "Lines, separators", "border"),
]

_HEX_RE = re.compile(r"#[0-9A-Fa-f]{6}\b")

_TONE_KEYS = (
    "serious_playful",
    "formal_casual",
    "respectful_irreverent",
    "matter_of_fact_enthusiastic",
)

_DESIGN_ELEMENT_KEYS = (
    "patterns_textures",
    "icon_style",
    "image_treatment",
    "decorative_elements",
)

# Keep in sync with frontend `brandUiIdentity.VISUAL_PARAGRAPH_SPLIT` (same pattern).
_MULTILINE_PARAGRAPH_SPLIT = re.compile(r"\n\s*\n+")


def _normalize_visual_prose_whitespace(text: str) -> str:
    """Turn stray single line breaks (e.g. legacy one-keyword-per-line) into spaces; keep paragraph breaks."""
    s = (text or "").replace("\r\n", "\n").strip()
    if not s:
        return ""
    blocks = _MULTILINE_PARAGRAPH_SPLIT.split(s)
    out: list[str] = []
    for block in blocks:
        inner = re.sub(r"\s+", " ", block.replace("\n", " ")).strip()
        if inner:
            out.append(inner)
    return "\n\n".join(out)


def _clip_visual_prose(raw: Any) -> str:
    """Short prose for visual_style; legacy keyword arrays joined into one string."""
    if isinstance(raw, list):
        parts = [str(x).strip() for x in raw if str(x).strip()]
        if not parts:
            return ""
        joined = " ".join(parts)
        return _normalize_visual_prose_whitespace(joined)[:_VISUAL_STYLE_TEXT_MAX]
    return _normalize_visual_prose_whitespace(str(raw or ""))[:_VISUAL_STYLE_TEXT_MAX]


def _normalize_visual_style(raw: Any) -> dict[str, Any]:
    """Prose summaries for visual identity; empty strings allowed."""
    empty_de = {k: "" for k in _DESIGN_ELEMENT_KEYS}
    defaults: dict[str, Any] = {
        "visual_mood_aesthetic": "",
        "style_guidelines": "",
        "design_elements": empty_de,
    }
    if not isinstance(raw, dict):
        return defaults

    out = dict(defaults)
    out["visual_mood_aesthetic"] = _clip_visual_prose(raw.get("visual_mood_aesthetic"))
    out["style_guidelines"] = _clip_visual_prose(raw.get("style_guidelines"))

    de_in = raw.get("design_elements")
    de_out = dict(empty_de)
    if isinstance(de_in, dict):
        for k in _DESIGN_ELEMENT_KEYS:
            de_out[k] = _clip_visual_prose(de_in.get(k))
    out["design_elements"] = de_out
    return out


def _empty_visual_style() -> dict[str, Any]:
    return _normalize_visual_style({})


def _normalize_tone_calibration(raw: Any) -> dict[str, int]:
    """0 = left pole, 100 = right pole, 50 = balanced."""
    defaults = {k: 50 for k in _TONE_KEYS}
    if not isinstance(raw, dict):
        return defaults
    out = dict(defaults)
    for k in _TONE_KEYS:
        try:
            v = int(raw.get(k))
            out[k] = max(0, min(100, v))
        except (TypeError, ValueError):
            pass
    return out


def _parse_first_hex(text: str) -> str | None:
    m = _HEX_RE.search(text)
    return m.group(0).upper() if m else None


def _normalize_hex(raw: str) -> str:
    s = (raw or "").strip()
    if not s:
        return "#CCCCCC"
    if s.startswith("#"):
        if len(s) == 4 and s[1:].isalnum():  # #RGB
            r, g, b = s[1], s[2], s[3]
            s = f"#{r}{r}{g}{g}{b}{b}"
        if len(s) == 7:
            return s.upper()
        m = _parse_first_hex(s)
        return m if m else "#CCCCCC"
    m = _parse_first_hex("#" + s)
    return m if m else "#CCCCCC"


def _ui_identity_fallback_from_guidelines(
    guidelines: dict[str, Any],
    brand_hint: str,
) -> dict[str, Any]:
    """Build ui_identity when pass3 is skipped (artifact mode) or failed."""
    primary = guidelines.get("primary_colors")
    fonts = guidelines.get("fonts")
    hexes: list[str] = []
    if isinstance(primary, list):
        for line in primary:
            hx = _parse_first_hex(str(line))
            if hx:
                hexes.append(hx)
            if len(hexes) >= 6:
                break
    while len(hexes) < 6:
        hexes.append("#CCCCCC")

    roles: list[dict[str, Any]] = []
    for i, (title, usage, surface) in enumerate(_DEFAULT_COLOR_ROLES):
        roles.append({
            "role": title,
            "usage": usage,
            "hex": hexes[i],
            "hierarchy_rank": i + 1,
            "surface": surface,
        })

    h1 = h2 = body = "Inter"
    if isinstance(fonts, list) and fonts:
        fam0 = fonts[0] if isinstance(fonts[0], dict) else {}
        if isinstance(fam0, dict) and fam0.get("family"):
            h1 = h2 = body = str(fam0["family"])
        if len(fonts) > 1 and isinstance(fonts[1], dict) and fonts[1].get("family"):
            body = str(fonts[1]["family"])

    return {
        "brand_name": brand_hint.strip() or "Brand",
        "color_roles": roles,
        "typography": {
            "heading_1": {"family": h1},
            "heading_2": {"family": h2},
            "body": {"family": body},
        },
        "tone_calibration": {k: 50 for k in _TONE_KEYS},
        "visual_style": _empty_visual_style(),
    }


def _normalize_ui_identity(
    raw: dict[str, Any] | None,
    merged_guidelines: dict[str, Any],
    brand_hint: str,
) -> dict[str, Any]:
    """Ensure ui_identity has roles, typography, tone_calibration, and visual_style."""
    base = _ui_identity_fallback_from_guidelines(merged_guidelines, brand_hint)
    if not isinstance(raw, dict):
        return base

    name = str(raw.get("brand_name") or "").strip() or base["brand_name"]
    typo = raw.get("typography")
    out_typo = dict(base["typography"])
    if isinstance(typo, dict):
        for key in ("heading_1", "heading_2", "body"):
            slot = typo.get(key)
            if isinstance(slot, dict) and slot.get("family"):
                fam = str(slot["family"]).strip()
                if fam:
                    entry: dict[str, Any] = {"family": fam[:120]}
                    if str(slot.get("source") or "") == "custom" and slot.get("custom_url"):
                        entry["source"] = "custom"
                        entry["custom_url"] = str(slot["custom_url"])[:2048]
                    out_typo[key] = entry

    roles_in = raw.get("color_roles")
    roles: list[dict[str, Any]] = []
    if isinstance(roles_in, list):
        for item in roles_in[:12]:
            if not isinstance(item, dict):
                continue
            role = str(item.get("role") or "").strip() or "Color"
            usage = str(item.get("usage") or "").strip() or "Usage"
            hx = _normalize_hex(str(item.get("hex") or ""))
            try:
                rank = int(item.get("hierarchy_rank"))
            except (TypeError, ValueError):
                rank = 3
            rank = max(1, min(6, rank))
            surface = str(item.get("surface") or "neutral").strip() or "neutral"
            roles.append({
                "role": role[:120],
                "usage": usage[:240],
                "hex": hx,
                "hierarchy_rank": rank,
                "surface": surface[:40],
            })

    if len(roles) < 6:
        fb = base["color_roles"]
        for i in range(len(roles), 6):
            roles.append(dict(fb[i]))
    elif len(roles) > 6:
        roles = roles[:6]

    # Sort by hierarchy_rank so UI can show most-used first
    roles.sort(key=lambda r: int(r.get("hierarchy_rank", 6)))

    return {
        "brand_name": name[:200],
        "color_roles": roles,
        "typography": out_typo,
        "tone_calibration": _normalize_tone_calibration(raw.get("tone_calibration")),
        "visual_style": _normalize_visual_style(raw.get("visual_style")),
    }


async def _run_brand_extract_passes(
    parts: list[dict[str, Any]],
    artifact_mode: bool,
) -> tuple[dict[str, Any], dict[str, Any], dict[str, Any]]:
    """Run visual + voice (+ ui_identity for full extraction) in parallel."""
    pass1_sys = PASS1_ARTIFACT_SYSTEM if artifact_mode else PASS1_SYSTEM
    pass2_sys = PASS2_ARTIFACT_SYSTEM if artifact_mode else PASS2_SYSTEM
    pass1_schema = _pass1_schema_for_artifact() if artifact_mode else PASS1_SCHEMA
    pass2_schema = _pass2_schema_for_artifact() if artifact_mode else PASS2_SCHEMA
    tok = ARTIFACT_MAX_OUTPUT_TOKENS if artifact_mode else None

    async def p1() -> dict[str, Any]:
        return await _gemini_extract_pass(
            parts=parts,
            system_instruction=pass1_sys,
            response_schema=pass1_schema,
            pass_name="visual",
            max_output_tokens=tok,
        )

    async def p2() -> dict[str, Any]:
        return await _gemini_extract_pass(
            parts=parts,
            system_instruction=pass2_sys,
            response_schema=pass2_schema,
            pass_name="voice",
            max_output_tokens=tok,
        )

    if artifact_mode:
        pass1_result, pass2_result = await asyncio.gather(p1(), p2())
        return pass1_result, pass2_result, {}

    async def p3() -> dict[str, Any]:
        return await _gemini_extract_pass(
            parts=parts,
            system_instruction=PASS3_SYSTEM,
            response_schema=PASS3_SCHEMA,
            pass_name="ui_identity",
            max_output_tokens=4096,
        )

    return await asyncio.gather(p1(), p2(), p3())


# ---------------------------------------------------------------------------
# Artifact mode — concise extraction for interactive HTML poll slides
# (avoids "exhaustive" deck dumps that overwhelm the artifact LLM)
# ---------------------------------------------------------------------------

PASS1_ARTIFACT_SYSTEM = (
    "You extract VISUAL design hints for ONE small interactive HTML poll slide (web UI), "
    "not a full brand manual or strategy deck.\n\n"
    "CRITICAL RULES:\n"
    "- Be SHORT. One line per list item where possible (hex + brief role, under ~100 chars per line).\n"
    "- Do NOT copy long paragraphs, roadmap bullets, governance text, or slide content from the document.\n"
    "- Colours: capture the main palette only (dominant backgrounds, text, accent). Skip exhaustive lists.\n"
    "- Fonts: at most 3 families; usage is a few words (e.g. 'titles', 'body').\n"
    "- For every string field: at most 2–3 sentences OR stay under the schema maxLength.\n"
    "- Ignore photography subjects, cityscapes, and narrative unless they change UI style in one phrase.\n"
    "- gradient_styles: at most 2 short entries.\n\n"
    "Return JSON matching the schema exactly."
)

PASS2_ARTIFACT_SYSTEM = (
    "You extract voice and HIGH-LEVEL visual direction for ONE small interactive HTML poll slide.\n\n"
    "CRITICAL RULES:\n"
    "- visual_style and tone_of_voice: each must stay under the schema maxLength — concise prose only.\n"
    "- key_principles: at most 5 items. Each item ONE short line about look, feel, or voice. "
    "Do NOT paste deck bullets, QEF roadmaps, objectives, or long quotes.\n"
    "- messaging_framework: at most 2 short phrases (taglines), not paragraphs.\n"
    "- dos_and_donts: at most 3 short design-related rules.\n"
    "- Skip business strategy, implementation plans, and document-specific content.\n\n"
    "Return JSON matching the schema exactly."
)

ARTIFACT_MAX_OUTPUT_TOKENS = 4096


def _pass1_schema_for_artifact() -> dict[str, Any]:
    """Tighter arrays and string caps so Gemini cannot return book-length fields."""
    s = copy.deepcopy(PASS1_SCHEMA)
    props = s["properties"]
    short_line = {"type": "string", "maxLength": 120}
    med_line = {"type": "string", "maxLength": 200}

    props["primary_colors"] = {
        "type": "array",
        "maxItems": 6,
        "items": short_line,
        "description": "Dominant colours with hex and brief usage (one line each)",
    }
    props["secondary_colors"] = {"type": "array", "maxItems": 4, "items": short_line}
    props["accent_colors"] = {"type": "array", "maxItems": 4, "items": short_line}
    props["gradient_styles"] = {"type": "array", "maxItems": 2, "items": med_line}
    props["fonts"] = {
        "type": "array",
        "maxItems": 3,
        "items": {
            "type": "object",
            "properties": {
                "family": {"type": "string", "maxLength": 80},
                "weights": {
                    "type": "array",
                    "maxItems": 4,
                    "items": {"type": "string", "maxLength": 24},
                },
                "usage": {"type": "string", "maxLength": 120},
            },
            "required": ["family"],
            "additionalProperties": False,
        },
    }
    props["typography_hierarchy"] = {"type": "string", "maxLength": 320}
    props["logo_description"] = {"type": "string", "maxLength": 180}
    props["logo_colors"] = {"type": "array", "maxItems": 6, "items": short_line}
    props["patterns_and_textures"] = {"type": "string", "maxLength": 200}
    props["brand_shapes"] = {"type": "string", "maxLength": 200}
    props["background_styles"] = {"type": "string", "maxLength": 260}
    props["iconography_style"] = {"type": "string", "maxLength": 200}
    props["illustration_style"] = {"type": "string", "maxLength": 160}
    props["photography_style"] = {"type": "string", "maxLength": 160}
    props["spacing_and_layout"] = {"type": "string", "maxLength": 260}
    props["animation_motion"] = {"type": "string", "maxLength": 120}
    return s


def _pass2_schema_for_artifact() -> dict[str, Any]:
    s = copy.deepcopy(PASS2_SCHEMA)
    props = s["properties"]
    props["visual_style"] = {
        "type": "string",
        "maxLength": 400,
        "description": "Brief visual mood for a web slide (not a full brand essay)",
    }
    props["tone_of_voice"] = {
        "type": "string",
        "maxLength": 320,
        "description": "Brief tone summary for UI copy",
    }
    props["key_principles"] = {
        "type": "array",
        "maxItems": 5,
        "items": {"type": "string", "maxLength": 140},
        "description": "At most 5 short lines; design/voice only, not deck strategy bullets",
    }
    props["messaging_framework"] = {
        "type": "array",
        "maxItems": 2,
        "items": {"type": "string", "maxLength": 100},
        "description": "Short taglines only",
    }
    props["dos_and_donts"] = {
        "type": "array",
        "maxItems": 3,
        "items": {"type": "string", "maxLength": 120},
        "description": "Short design-related rules only",
    }
    return s


# ---------------------------------------------------------------------------
# Image extraction from PDF / PPTX
# ---------------------------------------------------------------------------

def _extract_images_from_pdf(file_bytes: bytes) -> list[dict[str, Any]]:
    """Extract embedded images from a PDF as data URL candidates."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        logger.info("PyMuPDF not installed — skipping PDF image extraction")
        return []

    candidates: list[dict[str, Any]] = []
    seen_hashes: set[str] = set()

    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
    except Exception as exc:
        logger.warning("Failed to open PDF for image extraction: %s", exc)
        return []

    try:
        for page_num in range(min(len(doc), 30)):
            page = doc[page_num]
            for img_info in page.get_images(full=True):
                xref = img_info[0]
                try:
                    base_image = doc.extract_image(xref)
                except Exception:
                    continue
                if not base_image or not base_image.get("image"):
                    continue

                width = base_image.get("width", 0)
                height = base_image.get("height", 0)
                if width < IMAGE_MIN_DIMENSION or height < IMAGE_MIN_DIMENSION:
                    continue

                img_bytes = base_image["image"]
                img_hash = hashlib.md5(img_bytes).hexdigest()
                if img_hash in seen_hashes:
                    continue
                seen_hashes.add(img_hash)

                ext = base_image.get("ext", "png")
                mime = f"image/{ext}" if ext != "jpg" else "image/jpeg"
                b64 = base64.standard_b64encode(img_bytes).decode("ascii")
                if len(b64) > IMAGE_MAX_DATA_URL_BYTES:
                    continue

                candidates.append({
                    "data_url": f"data:{mime};base64,{b64}",
                    "width": width,
                    "height": height,
                    "page": page_num + 1,
                })
    finally:
        doc.close()

    candidates.sort(key=lambda c: c["width"] * c["height"], reverse=True)
    return candidates[:IMAGE_MAX_CANDIDATES]


def _extract_images_from_pptx(file_bytes: bytes) -> list[dict[str, Any]]:
    """Extract embedded images from a PPTX file."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        logger.info("python-pptx not installed — skipping PPTX image extraction")
        return []

    candidates: list[dict[str, Any]] = []
    seen_hashes: set[str] = set()

    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except Exception as exc:
        logger.warning("Failed to open PPTX for image extraction: %s", exc)
        return []

    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_blob = shape.image.blob
                    content_type = shape.image.content_type or "image/png"
                except Exception:
                    continue

                img_hash = hashlib.md5(img_blob).hexdigest()
                if img_hash in seen_hashes:
                    continue
                seen_hashes.add(img_hash)

                b64 = base64.standard_b64encode(img_blob).decode("ascii")
                if len(b64) > IMAGE_MAX_DATA_URL_BYTES:
                    continue

                width = int(shape.width / 9525) if shape.width else 0
                height = int(shape.height / 9525) if shape.height else 0
                if width < IMAGE_MIN_DIMENSION or height < IMAGE_MIN_DIMENSION:
                    continue

                candidates.append({
                    "data_url": f"data:{content_type};base64,{b64}",
                    "width": width,
                    "height": height,
                    "page": slide_num,
                })

    candidates.sort(key=lambda c: c["width"] * c["height"], reverse=True)
    return candidates[:IMAGE_MAX_CANDIDATES]


def _extract_images_from_upload(
    file_bytes: bytes, content_type: str,
) -> list[dict[str, Any]]:
    """Extract candidate logo/brand images from an uploaded file."""
    if content_type in SUPPORTED_IMAGE_TYPES:
        b64 = base64.standard_b64encode(file_bytes).decode("ascii")
        if len(b64) <= IMAGE_MAX_DATA_URL_BYTES:
            return [{"data_url": f"data:{content_type};base64,{b64}", "width": 0, "height": 0, "page": 1}]
        return []
    if content_type == "application/pdf":
        return _extract_images_from_pdf(file_bytes)
    if content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return _extract_images_from_pptx(file_bytes)
    return []


# ---------------------------------------------------------------------------
# Gemini helpers
# ---------------------------------------------------------------------------

def _gemini_base_url() -> str:
    return settings.gemini_base_url.rstrip("/")


def _gemini_generate_url(model: str) -> str:
    return f"{_gemini_base_url()}/models/{model}:generateContent"


def _gemini_upload_url() -> str:
    return f"{_gemini_base_url().replace('/v1beta', '')}/upload/v1beta/files"


async def _upload_to_gemini_files_api(
    file_bytes: bytes,
    content_type: str,
    filename: str,
    api_key: str,
) -> str:
    """Upload a file to the Gemini File API and return its URI."""
    boundary = f"----GeminiBoundary{uuid.uuid4().hex}"
    metadata = f'{{"file": {{"displayName": "{filename}"}}}}'

    body = (
        f"--{boundary}\r\n"
        f"Content-Type: application/json; charset=UTF-8\r\n\r\n"
        f"{metadata}\r\n"
        f"--{boundary}\r\n"
        f"Content-Type: {content_type}\r\n\r\n"
    ).encode() + file_bytes + f"\r\n--{boundary}--".encode()

    try:
        async with httpx.AsyncClient(
            timeout=httpx.Timeout(120.0, connect=10.0)
        ) as client:
            response = await client.post(
                _gemini_upload_url(),
                params={"key": api_key, "uploadType": "multipart"},
                content=body,
                headers={"Content-Type": f"multipart/related; boundary={boundary}"},
            )
    except httpx.TimeoutException as exc:
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail="File upload to Gemini timed out",
        ) from exc
    except httpx.RequestError as exc:
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail="File upload to Gemini failed",
        ) from exc

    if response.status_code >= 400:
        detail = ""
        try:
            detail = response.json().get("error", {}).get("message", "")
        except Exception:
            pass
        logger.error("Gemini file upload error: status=%d detail=%s", response.status_code, detail)
        raise HTTPException(
            status_code=status.HTTP_502_BAD_GATEWAY,
            detail=f"File upload failed: {detail or response.status_code}",
        )

    try:
        file_uri = response.json()["file"]["uri"]
    except (KeyError, TypeError) as exc:
        logger.error("Gemini file upload: unexpected response shape")
        raise HTTPException(
            status_code=status.HTTP_502_BAD_GATEWAY,
            detail="File upload returned an unexpected response",
        ) from exc

    return file_uri


async def _delete_gemini_file(file_uri: str, api_key: str) -> None:
    """Best-effort deletion of an uploaded Gemini file."""
    try:
        file_name = file_uri.rstrip("/").rsplit("/", 1)[-1]
        async with httpx.AsyncClient(timeout=10.0) as client:
            await client.delete(
                f"{_gemini_base_url()}/files/{file_name}",
                params={"key": api_key},
            )
    except Exception as exc:
        logger.warning("Could not delete Gemini file %s: %s", file_uri, exc)


async def _gemini_extract_pass(
    *,
    parts: list[dict[str, Any]],
    system_instruction: str,
    response_schema: dict[str, Any],
    pass_name: str,
    max_output_tokens: int | None = None,
) -> dict[str, Any]:
    """Run a single schema-enforced Gemini extraction pass.

    Uses `responseJsonSchema` to force Gemini to populate every field
    in the schema, preventing the raw_notes dumping problem.
    Falls back to non-schema mode if the schema exceeds Gemini's state limit.
    """
    api_key = settings.gemini_api_key
    if not api_key:
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail="Gemini API key is not configured",
        )

    model = settings.gemini_brand_extract_model or settings.gemini_model
    url = _gemini_generate_url(model)

    body: dict[str, Any] = {
        "contents": [{"parts": parts}],
        "systemInstruction": {"parts": [{"text": system_instruction}]},
        "generationConfig": {
            "temperature": 0.2,
            "maxOutputTokens": max_output_tokens if max_output_tokens is not None else 12000,
            "responseMimeType": "application/json",
            "responseJsonSchema": response_schema,
        },
    }

    result = await _call_gemini(url, api_key, body, pass_name)
    if result is not None:
        return result

    # Fallback: retry without schema if schema caused a state overflow error
    logger.warning("Pass %s: retrying without responseJsonSchema", pass_name)
    body["generationConfig"].pop("responseJsonSchema", None)
    result = await _call_gemini(url, api_key, body, pass_name)
    if result is not None:
        return result

    return {}


async def _call_gemini(
    url: str,
    api_key: str,
    body: dict[str, Any],
    pass_name: str,
) -> dict[str, Any] | None:
    """Execute a Gemini generateContent call. Returns parsed JSON or None on failure."""
    try:
        async with httpx.AsyncClient(
            timeout=httpx.Timeout(EXTRACT_TIMEOUT_SECONDS, connect=10.0)
        ) as client:
            response = await client.post(url, params={"key": api_key}, json=body)
    except httpx.TimeoutException:
        logger.error("Gemini brand extract timeout: pass=%s", pass_name)
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail=f"Brand extraction timed out ({pass_name})",
        )
    except httpx.RequestError as exc:
        logger.error("Gemini brand extract request error: pass=%s error=%s", pass_name, exc)
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail=f"Brand extraction request failed ({pass_name})",
        )

    if response.status_code >= 400:
        detail = ""
        try:
            detail = response.json().get("error", {}).get("message", "")
        except Exception:
            pass
        # If schema caused a "too many states" error, return None so caller can retry
        if "too many" in detail.lower() or "state" in detail.lower():
            logger.warning("Gemini schema overflow in pass %s: %s", pass_name, detail)
            return None
        logger.error("Gemini brand extract error: pass=%s status=%d detail=%s", pass_name, response.status_code, detail)
        raise HTTPException(
            status_code=status.HTTP_502_BAD_GATEWAY,
            detail=f"Brand extraction failed ({pass_name}): {detail or response.status_code}",
        )

    try:
        data = response.json()
        text = data["candidates"][0]["content"]["parts"][0]["text"]
    except (KeyError, IndexError, TypeError):
        logger.error("Gemini brand extract: unexpected response shape in pass %s", pass_name)
        return None

    try:
        return _json.loads(text)
    except _json.JSONDecodeError:
        logger.warning("Gemini brand extract: invalid JSON in pass %s", pass_name)
        return None


# ---------------------------------------------------------------------------
# Main endpoint
# ---------------------------------------------------------------------------

@router.post("/extract")
async def extract_brand_profile(
    user: AuthUser = Depends(get_library_user),
    file: UploadFile | None = File(default=None),
    url: str | None = Form(default=None),
    purpose: str = Form(default="full"),
) -> dict[str, Any]:
    """Extract brand guidelines from an uploaded file or a website URL.

    Uses a parallel extraction strategy:
      Pass 1 — Visual Identity (colours, fonts, patterns, asset styles)
      Pass 2 — Brand Voice & Strategy (tone, principles, messaging, dos/don'ts)
      Pass 3 — UI identity (full purpose only): brand name, 6 color roles, typography H1/H2/body

    Passes run concurrently with schema enforcement so Gemini is forced
    to populate every field instead of dumping into raw_notes.

    `purpose`:
      - `full` (default): exhaustive extraction for saved brand profiles / reference.
      - `artifact`: concise caps and prompts for interactive poll HTML (avoid deck dumps).
    """

    purpose_norm = (purpose or "full").strip().lower()
    if purpose_norm not in ("full", "artifact"):
        purpose_norm = "full"
    artifact_mode = purpose_norm == "artifact"

    if file and url:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Provide either a file or a URL, not both",
        )

    if not file and not url:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Provide a file upload or a URL",
        )

    # --- Build the content parts (shared by both passes) ---
    parts: list[dict[str, Any]] = []
    extracted_images: list[dict[str, Any]] = []

    if file:
        content_type = file.content_type or "application/octet-stream"
        if content_type not in SUPPORTED_UPLOAD_TYPES:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Unsupported file type: {content_type}. "
                "Supported: PDF, PPTX, PNG, JPEG, GIF, WEBP",
            )

        file_bytes = await file.read()
        file_size = len(file_bytes)

        if file_size > EXTRACT_MAX_FILE_SIZE:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"File too large (max {EXTRACT_MAX_FILE_SIZE // (1024 * 1024)} MB)",
            )

        api_key = settings.gemini_api_key
        if not api_key:
            raise HTTPException(
                status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
                detail="Gemini API key is not configured",
            )

        # Extract embedded images (logos, assets) — runs synchronously, fast
        extracted_images = _extract_images_from_upload(file_bytes, content_type)

        file_uri: str | None = None
        try:
            if file_size > INLINE_MAX_FILE_SIZE:
                logger.info(
                    "Uploading large file (%d MB) via Gemini File API",
                    file_size // (1024 * 1024),
                )
                file_uri = await _upload_to_gemini_files_api(
                    file_bytes, content_type, file.filename or "file", api_key
                )
                parts.append({
                    "fileData": {"mimeType": content_type, "fileUri": file_uri}
                })
            else:
                b64 = base64.standard_b64encode(file_bytes).decode("ascii")
                parts.append({
                    "inlineData": {"mimeType": content_type, "data": b64}
                })

            parts.append({
                "text": f"Extract brand guidelines from this uploaded file ({file.filename or 'file'})."
            })

            source_type = "file"
            source_filename = file.filename or ""

            pass1_result, pass2_result, pass3_result = await _run_brand_extract_passes(
                parts, artifact_mode
            )
        finally:
            if file_uri:
                await _delete_gemini_file(file_uri, api_key)

    else:
        assert url is not None
        page_snippet = ""
        try:
            async with httpx.AsyncClient(
                timeout=httpx.Timeout(15.0, connect=5.0),
                follow_redirects=True,
                max_redirects=5,
            ) as client:
                page_response = await client.get(url, headers={
                    "User-Agent": "Mozilla/5.0 (compatible; PrezoBot/1.0)",
                    "Accept": "text/html,application/xhtml+xml",
                })
                if page_response.status_code < 400:
                    raw_html = page_response.text[:60000]
                    raw_html = re.sub(
                        r"<(script|style)[^>]*>[\s\S]*?</\1>",
                        "",
                        raw_html,
                        flags=re.IGNORECASE,
                    )
                    page_snippet = raw_html[:30000]
        except Exception as exc:
            logger.info("Could not fetch URL %s for brand extraction: %s", url, exc)

        if page_snippet:
            parts.append({
                "text": (
                    f"Analyse the brand identity of the website at {url}.\n"
                    "Below is the page HTML. Extract colours, fonts, visual style, "
                    "and brand guidelines from the actual content.\n\n"
                    f"{page_snippet}"
                )
            })
        else:
            parts.append({
                "text": (
                    f"Analyse the brand identity of the website at: {url}\n"
                    "Based on the URL, describe the likely visual identity, colours, fonts, "
                    "and brand style. Extract brand guidelines as best you can."
                )
            })

        source_type = "url"
        source_filename = url

        pass1_result, pass2_result, pass3_result = await _run_brand_extract_passes(
            parts, artifact_mode
        )

    # --- Merge pass results into a single guidelines object ---
    guidelines: dict[str, Any] = {}
    if isinstance(pass1_result, dict):
        guidelines.update(pass1_result)
    if isinstance(pass2_result, dict):
        guidelines.update(pass2_result)

    brand_hint = ""
    if file:
        fn = source_filename or ""
        base = fn.rsplit("/", 1)[-1].rsplit("\\", 1)[-1]
        brand_hint = base.rsplit(".", 1)[0] if base else ""
    else:
        try:
            from urllib.parse import urlparse

            host = urlparse(url or "").hostname or ""
            brand_hint = host.split(".")[0] if host else ""
        except Exception:
            brand_hint = ""

    ui_raw = pass3_result if isinstance(pass3_result, dict) else {}
    guidelines["ui_identity"] = _normalize_ui_identity(
        ui_raw if not artifact_mode else None,
        guidelines,
        brand_hint,
    )

    result: dict[str, Any] = {
        "source_type": source_type,
        "source_filename": source_filename,
        "guidelines": guidelines,
        "raw_summary": "",
        "extraction_purpose": purpose_norm,
    }

    if extracted_images:
        result["extracted_images"] = extracted_images

    return result
